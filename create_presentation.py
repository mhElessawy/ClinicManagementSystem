#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Clinic Management System - PowerPoint Presentation Generator
Creates a professional bilingual presentation in Arabic and English
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create a professional bilingual PowerPoint presentation"""

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    DARK_BLUE = RGBColor(44, 62, 80)      # #2c3e50
    LIGHT_BLUE = RGBColor(52, 152, 219)   # #3498db
    RED_ACCENT = RGBColor(231, 76, 60)    # #e74c3c
    GREEN = RGBColor(39, 174, 96)         # #27ae60
    ORANGE = RGBColor(243, 156, 18)       # #f39c12
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(236, 240, 241)  # #ecf0f1

    # ============ Slide 1: Title Slide ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "نظام إدارة العيادات"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Clinic Management System"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = "حل احترافي لإدارة المنشآت الصحية\nProfessional Healthcare Management Solution"
    p = desc_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Technology
    tech_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.5), Inches(5), Inches(0.5))
    tech_frame = tech_box.text_frame
    tech_frame.text = "Powered by ASP.NET Core 8 MVC"
    p = tech_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # ============ Slide 2: Overview ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "نظرة عامة | Overview"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Arabic content
    p = tf.paragraphs[0]
    p.text = "نظام إدارة العيادات هو حل برمجي متكامل وحديث مصمم لإدارة العيادات الطبية بكفاءة عالية."
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT
    p.space_after = Pt(12)

    # English content
    p = tf.add_paragraph()
    p.text = "Clinic Management System is a comprehensive and modern software solution designed for efficient medical clinic management."
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(85, 85, 85)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(18)

    # Features
    p = tf.add_paragraph()
    p.text = "✓ يوفر النظام إدارة شاملة للأطباء، المرضى، المواعيد، والتشخيصات الطبية"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.RIGHT
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "✓ Complete management of doctors, patients, appointments, and medical diagnoses"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(18)

    p = tf.add_paragraph()
    p.text = "✓ مبني على أحدث تقنيات Microsoft مع تركيز على الأمان والأداء"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.RIGHT
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "✓ Built on the latest Microsoft technologies with focus on security and performance"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.LEFT

    # Statistics boxes
    stats = [
        ("10", "Database Tables\nجداول قاعدة البيانات", 1.2),
        ("16", "Controllers\nوحدات التحكم", 3.2),
        ("50+", "Views\nصفحات العرض", 5.2),
        ("5,000+", "Lines of Code\nسطر برمجي", 7.2)
    ]

    for num, label, left in stats:
        # Box
        box = slide.shapes.add_shape(1, Inches(left), Inches(5.3), Inches(1.6), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = DARK_BLUE

        # Number
        num_box = slide.shapes.add_textbox(Inches(left), Inches(5.4), Inches(1.6), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(left), Inches(6.1), Inches(1.6), Inches(0.7))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 3: Core Features ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "الميزات الأساسية | Core Features"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Feature cards
    features = [
        ("إدارة الأطباء", "Doctor Management", "Complete profiles, specialties, credentials", 1, 1.3),
        ("إدارة المرضى", "Patient Management", "Patient records, history, doctor assignment", 3.8, 1.3),
        ("نظام المواعيد", "Appointments System", "Schedule, track, manage appointments", 6.6, 1.3),
        ("التشخيصات الطبية", "Medical Diagnoses", "Diagnosis documentation with attachments", 1, 3.8),
        ("إدارة المساعدين", "Staff Management", "Doctor assistants with login credentials", 3.8, 3.8),
        ("التقارير المتقدمة", "Advanced Reports", "Excel exports with filters", 6.6, 3.8)
    ]

    for ar_title, en_title, desc, left, top in features:
        # Card background
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.4), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BLUE
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(2)

        # Arabic title
        ar_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.2), Inches(2.2), Inches(0.5))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English title
        en_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.7), Inches(2.2), Inches(0.4))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_title
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.2), Inches(2.2), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 4: User Roles ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "أدوار المستخدمين | User Roles"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Roles
    roles = [
        ("المدير الأعلى", "Super Admin", "Full system access", DARK_BLUE, 1.5, 1.5),
        ("المدير", "Admin", "Manage departments, doctors, patients", LIGHT_BLUE, 5.5, 1.5),
        ("الطبيب", "Doctor", "Own patients & diagnoses", GREEN, 1.5, 3.5),
        ("المساعد", "Assistant", "Limited access, tied to doctor", ORANGE, 5.5, 3.5),
        ("موظف الاستقبال", "Receptionist", "Patient management only", RED_ACCENT, 3.5, 5.5)
    ]

    for ar_role, en_role, desc, color, left, top in roles:
        # Role card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(2)

        # Arabic role
        ar_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(2.6), Inches(0.4))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_role
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English role
        en_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.6), Inches(2.6), Inches(0.3))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_role
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1), Inches(2.6), Inches(0.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 5: Technology Stack ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "التقنيات المستخدمة | Technology Stack"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Backend section
    backend_title = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(0.5))
    tf = backend_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Backend Technologies | تقنيات الخادم"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Backend technologies
    backend_techs = [
        ("ASP.NET Core 8.0", 1.5, 2),
        ("Entity Framework Core 8.0", 3.7, 2),
        ("SQL Server 2019+", 5.9, 2),
        ("C# 12", 1.5, 3),
        ("BCrypt.Net-Next", 3.7, 3)
    ]

    for tech, left, top in backend_techs:
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.color.rgb = LIGHT_BLUE
        box.line.width = Pt(2)

        text_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(1.8), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Frontend section
    frontend_title = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    tf = frontend_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Frontend Technologies | تقنيات الواجهة"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Frontend technologies
    frontend_techs = [
        ("Bootstrap 5.3.0", 1.5, 4.7),
        ("Font Awesome 6.4.0", 3.7, 4.7),
        ("jQuery", 5.9, 4.7),
        ("Responsive Design", 1.5, 5.7),
        ("ClosedXML (Excel)", 3.7, 5.7),
        ("RTL Support", 5.9, 5.7)
    ]

    for tech, left, top in frontend_techs:
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = GREEN
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(2)

        text_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(1.8), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ============ Slide 6: Database Structure ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "هيكل قاعدة البيانات | Database Structure"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Database tables
    tables = [
        ("UserInfo & Role", "إدارة المستخدمين والصلاحيات", "User authentication and permissions", 1.2, 1.3),
        ("Department & Specialist", "الأقسام والتخصصات", "Medical departments and specialties", 5.8, 1.3),
        ("DoctorInfo", "بيانات الأطباء", "Doctor profiles and credentials", 1.2, 2.8),
        ("Patient", "بيانات المرضى", "Patient records and information", 5.8, 2.8),
        ("PatientDiagnosis", "التشخيصات الطبية", "Medical diagnosis documentation", 1.2, 4.3),
        ("Appointment", "المواعيد", "Appointment scheduling and tracking", 5.8, 4.3),
        ("DoctorAssist", "مساعدو الأطباء", "Doctor assistants management", 1.2, 5.8),
        ("DoctorSubscription", "اشتراكات الأطباء", "Subscription management", 5.8, 5.8)
    ]

    for en_name, ar_desc, en_desc, left, top in tables:
        # Table card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.4), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = LIGHT_BLUE
        card.line.width = Pt(3)

        # Table name
        name_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(3.2), Inches(0.3))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.LEFT

        # Arabic description
        ar_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.45), Inches(3.2), Inches(0.3))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.RIGHT

        # English description
        en_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.75), Inches(3.2), Inches(0.3))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(85, 85, 85)
        p.alignment = PP_ALIGN.LEFT

    # ============ Slide 7: System Statistics ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "إحصائيات النظام | System Statistics"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Statistics
    stats_data = [
        ("16", "Controllers\nوحدات التحكم", LIGHT_BLUE, 1, 1.5),
        ("10", "Models\nالنماذج", GREEN, 3.5, 1.5),
        ("50+", "Views\nصفحات العرض", ORANGE, 6, 1.5),
        ("5", "User Roles\nأدوار المستخدمين", RED_ACCENT, 1, 3.5),
        ("40+", "Features\nالميزات", DARK_BLUE, 3.5, 3.5),
        ("4", "Report Types\nأنواع التقارير", LIGHT_BLUE, 6, 3.5),
        ("5,000+", "Lines of Code\nسطر برمجي", GREEN, 2.25, 5.5),
        ("25+", "Documentation Files\nملفات التوثيق", ORANGE, 4.75, 5.5)
    ]

    for num, label, color, left, top in stats_data:
        # Stat box
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.2), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(2)

        # Number
        num_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.2), Inches(2), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.85), Inches(2), Inches(0.5))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 8: Security Features ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "ميزات الأمان | Security Features"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Security features
    security_features = [
        ("تشفير كلمات المرور", "Password Encryption", "BCrypt hashing", 1.2, 1.5),
        ("إدارة الجلسات", "Session Management", "30-min timeout, HTTPOnly cookies", 5.4, 1.5),
        ("التحكم في الوصول", "Access Control", "Role-based permissions", 1.2, 3.3),
        ("استعلامات آمنة", "Secure Queries", "Parameterized via EF Core", 5.4, 3.3),
        ("حماية CSRF", "CSRF Protection", "Anti-Forgery tokens", 1.2, 5.1),
        ("سجل التدقيق", "Audit Trail", "Track all user actions", 5.4, 5.1)
    ]

    for ar_title, en_title, desc, left, top in security_features:
        # Feature card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.4), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = GREEN
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(2)

        # Arabic title
        ar_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(3), Inches(0.4))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English title
        en_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(3), Inches(0.35))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_title
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.95), Inches(3), Inches(0.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 9: Reports & Analytics ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "التقارير والتحليلات | Reports & Analytics"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Report types
    reports = [
        ("تقرير الأطباء", "Doctors Report", "Excel export with specialty filtering", LIGHT_BLUE, 1.5, 1.5),
        ("تقرير المرضى", "Patients Report", "Excel export with doctor filtering", GREEN, 5.5, 1.5),
        ("تقرير التشخيصات", "Diagnoses Report", "Filterable by date, doctor, patient", ORANGE, 1.5, 3.8),
        ("تقرير الإحصائيات", "Statistics Report", "System-wide analytics and insights", RED_ACCENT, 5.5, 3.8)
    ]

    for ar_title, en_title, desc, color, left, top in reports:
        # Report card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.2), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(3)

        # Arabic title
        ar_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(2.8), Inches(0.4))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English title
        en_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.65), Inches(2.8), Inches(0.35))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_title
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.1), Inches(2.8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Features
    features_box = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(1))
    tf = features_box.text_frame

    p = tf.paragraphs[0]
    p.text = "✓ فلترة متقدمة  ✓ تصدير Excel  ✓ نطاقات التاريخ  ✓ تقارير مخصصة"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "✓ Advanced Filtering  ✓ Excel Export  ✓ Date Ranges  ✓ Custom Reports"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(85, 85, 85)
    p.alignment = PP_ALIGN.CENTER

    # ============ Slide 10: Dashboard ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "لوحة التحكم | Dashboard"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Dashboard description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "لوحة تحكم ذكية مخصصة لكل دور مع عرض إحصائيات النظام في الوقت الفعلي"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Smart dashboard customized for each role with real-time system statistics"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(85, 85, 85)
    p.alignment = PP_ALIGN.CENTER

    # Dashboard features
    features_list = [
        ("📊 عرض إحصائيات النظام في الوقت الفعلي", "Real-time system statistics display"),
        ("📅 مواعيد اليوم والغد مباشرة", "Today's and tomorrow's appointments at a glance"),
        ("🔔 تنبيهات الاشتراكات والتجديد", "Subscription and renewal alerts"),
        ("👥 عدد الأطباء والمرضى والتشخيصات", "Count of doctors, patients, and diagnoses")
    ]

    top_position = 3
    for ar_feature, en_feature in features_list:
        # Feature box
        box = slide.shapes.add_shape(1, Inches(1.5), Inches(top_position), Inches(7), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = LIGHT_BLUE
        box.line.width = Pt(2)

        # Arabic text
        ar_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_position + 0.1), Inches(6.6), Inches(0.3))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_feature
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.RIGHT

        # English text
        en_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_position + 0.45), Inches(6.6), Inches(0.25))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_feature
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(85, 85, 85)
        p.alignment = PP_ALIGN.LEFT

        top_position += 1

    # ============ Slide 11: Use Cases ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "حالات الاستخدام | Use Cases"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Use cases
    use_cases = [
        ("🏥 العيادات الطبية", "Medical Clinics", "Complete management for private clinics", LIGHT_BLUE, 1.5, 1.5),
        ("🎓 التعليم", "Education", "Learn ASP.NET Core MVC and EF", GREEN, 5.5, 1.5),
        ("🏢 المراكز الصحية", "Health Centers", "Manage multiple departments", ORANGE, 1.5, 4),
        ("🚀 الشركات الناشئة", "Startups", "Foundation for healthcare platforms", RED_ACCENT, 5.5, 4)
    ]

    for ar_title, en_title, desc, color, left, top in use_cases:
        # Use case card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.2), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(3)

        # Arabic title
        ar_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(2.8), Inches(0.5))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English title
        en_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.85), Inches(2.8), Inches(0.4))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_title
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.3), Inches(2.8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Summary
    summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(7), Inches(0.8))
    tf = summary_box.text_frame

    p = tf.paragraphs[0]
    p.text = "حل مرن يناسب جميع أحجام المنشآت الطبية"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Flexible solution suitable for all sizes of medical facilities"
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(85, 85, 85)
    p.alignment = PP_ALIGN.CENTER

    # ============ Slide 12: Future Enhancements ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "التحسينات المستقبلية | Future Enhancements"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Future enhancements
    enhancements = [
        ("📱 تطبيق موبايل", "Mobile App", "iOS and Android"),
        ("✉️ إشعارات SMS/Email", "SMS/Email Notifications", "Automated reminders"),
        ("💊 وحدة الصيدلية", "Pharmacy Module", "Prescription management"),
        ("🧪 تكامل المختبر", "Laboratory Integration", "Lab test results"),
        ("💰 الفواتير والمدفوعات", "Billing & Payments", "Payment processing"),
        ("📊 تحليلات متقدمة", "Advanced Analytics", "AI-powered insights")
    ]

    top_position = 1.5
    for icon_ar, icon_en, desc in enhancements:
        # Enhancement box
        box = slide.shapes.add_shape(1, Inches(1.5), Inches(top_position), Inches(7), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(2)

        # Content
        content_box = slide.shapes.add_textbox(Inches(1.7), Inches(top_position + 0.1), Inches(6.6), Inches(0.55))
        tf = content_box.text_frame

        p = tf.paragraphs[0]
        p.text = f"{icon_ar}  |  {icon_en}  •  {desc}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        top_position += 0.9

    # ============ Slide 13: Key Benefits ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "الفوائد الرئيسية | Key Benefits"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Benefits
    benefits = [
        ("⚡ أداء عالي", "High Performance", "ASP.NET Core 8", LIGHT_BLUE, 1.2, 1.5),
        ("🛡️ آمن وموثوق", "Secure & Reliable", "Enterprise security", GREEN, 5.4, 1.5),
        ("📈 قابل للتوسع", "Scalable", "Grow with your business", ORANGE, 1.2, 3.3),
        ("🎨 واجهة عصرية", "Modern UI", "Bootstrap 5 responsive", RED_ACCENT, 5.4, 3.3),
        ("⚙️ سهل التخصيص", "Easy to Customize", "Modular architecture", DARK_BLUE, 1.2, 5.1),
        ("📚 موثق بالكامل", "Fully Documented", "25+ documentation files", LIGHT_BLUE, 5.4, 5.1)
    ]

    for ar_benefit, en_benefit, desc, color, left, top in benefits:
        # Benefit card
        card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.4), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = DARK_BLUE
        card.line.width = Pt(2)

        # Arabic benefit
        ar_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(3), Inches(0.4))
        tf = ar_box.text_frame
        p = tf.paragraphs[0]
        p.text = ar_benefit
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # English benefit
        en_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.6), Inches(3), Inches(0.35))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.text = en_benefit
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1), Inches(3), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ============ Slide 14: Conclusion ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "الخلاصة | Conclusion"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Main message
    msg_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(2.5))
    tf = msg_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "نظام إدارة عيادات متكامل وجاهز للإنتاج"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "Complete Production-Ready Clinic Management System"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)

    p = tf.add_paragraph()
    p.text = "حل برمجي شامل يجمع بين الأداء العالي، الأمان، وسهولة الاستخدام"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Comprehensive solution combining high performance, security, and ease of use"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Key points
    points_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(1))
    tf = points_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⭐ Modern Tech  •  🛡️ Secure  •  📈 Scalable  •  👥 User-Friendly"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # ============ Slide 15: Thank You ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = thanks_box.text_frame

    p = tf.paragraphs[0]
    p.text = "شكراً"
    p.font.size = Pt(70)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(2))
    tf = contact_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Clinic Management System"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "نظام إدارة العيادات"
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)

    p = tf.add_paragraph()
    p.text = "📧 Contact Us for More Information"
    p.font.size = Pt(22)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "🚀 Ready to Transform Your Clinic"
    p.font.size = Pt(20)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "/home/user/ClinicManagementSystem/ClinicManagementSystem_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        output_file = create_presentation()
        print(f"\n✓ PowerPoint presentation created successfully!")
        print(f"✓ File: {output_file}")
        print(f"✓ Total slides: 15")
        print(f"✓ Language: Bilingual (Arabic/English)")
    except Exception as e:
        print(f"✗ Error creating presentation: {str(e)}")
        raise
