#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Clinic Management System - Detailed PowerPoint Presentation with Website Pages
Creates a comprehensive bilingual presentation showing all website pages
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_detailed_presentation():
    """Create a detailed PowerPoint presentation with all website pages"""

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    DARK_BLUE = RGBColor(44, 62, 80)
    LIGHT_BLUE = RGBColor(52, 152, 219)
    RED_ACCENT = RGBColor(231, 76, 60)
    GREEN = RGBColor(39, 174, 96)
    ORANGE = RGBColor(243, 156, 18)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(236, 240, 241)
    PURPLE = RGBColor(142, 68, 173)

    def add_title_slide(title_ar, title_en, subtitle_ar="", subtitle_en=""):
        """Helper function to add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_ar
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = title_en
        p.font.size = Pt(44)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

        if subtitle_ar:
            subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{subtitle_ar}\n{subtitle_en}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    def add_page_slide(title_ar, title_en, features, components, color=LIGHT_BLUE):
        """Helper function to add a page description slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title_ar} | {title_en}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Features section
        features_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.4), Inches(3))
        tf = features_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "الميزات الرئيسية | Main Features:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

        for feature_ar, feature_en in features:
            p = tf.add_paragraph()
            p.text = f"• {feature_ar}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_BLUE
            p.space_after = Pt(4)
            p.level = 0

            p = tf.add_paragraph()
            p.text = feature_en
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(85, 85, 85)
            p.space_after = Pt(8)
            p.level = 1

        # Components section
        comp_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.4), Inches(2.8))
        tf = comp_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "المكونات | Components:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8)

        for comp in components:
            p = tf.add_paragraph()
            p.text = f"✓ {comp}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_BLUE
            p.space_after = Pt(4)

    # ============ Slide 1: Cover ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "نظام إدارة العيادات"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Clinic Management System"
    p.font.size = Pt(48)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "عرض شامل لجميع صفحات الموقع"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Complete Overview of All Website Pages"
    p.font.size = Pt(24)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    tech_box = slide.shapes.add_textbox(Inches(2.5), Inches(6), Inches(5), Inches(0.6))
    tf = tech_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ASP.NET Core 8 MVC • Bootstrap 5 • SQL Server"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # ============ Section: Authentication ============
    add_title_slide(
        "قسم المصادقة والأمان",
        "Authentication & Security Section",
        "صفحات تسجيل الدخول والأمان",
        "Login and Security Pages"
    )

    # ============ Page: Login ============
    add_page_slide(
        "صفحة تسجيل الدخول",
        "Login Page",
        [
            ("تصميم عصري بخلفية متدرجة جذابة", "Modern design with attractive gradient background"),
            ("اختيار نوع المستخدم (مدير، طبيب، مساعد)", "User type selection (Admin, Doctor, Assistant)"),
            ("أيقونات تفاعلية لكل نوع مستخدم", "Interactive icons for each user type"),
            ("نظام حماية CSRF", "CSRF protection system"),
            ("رسائل خطأ ونجاح واضحة", "Clear error and success messages")
        ],
        [
            "🎨 خلفية متدرجة (Gradient Background) باللون الأزرق والبنفسجي",
            "📋 نموذج تسجيل دخول (Login Form) مع اختيار نوع المستخدم",
            "👤 أيقونات Font Awesome للمستخدمين (Admin, Doctor, Assistant)",
            "🔒 حقول Username و Password مع Validation",
            "⚠️ Alert Boxes لعرض الرسائل (Success/Error)",
            "🎯 زر Login بتأثيرات Hover جذابة"
        ],
        PURPLE
    )

    # ============ Section: Dashboard ============
    add_title_slide(
        "لوحة التحكم الرئيسية",
        "Main Dashboard",
        "الصفحة الرئيسية للنظام",
        "System Home Page"
    )

    # ============ Page: Dashboard ============
    add_page_slide(
        "لوحة التحكم (Dashboard)",
        "Dashboard Page",
        [
            ("عرض ترحيبي باسم المستخدم ونوعه", "Welcome message with user name and type"),
            ("4 بطاقات إحصائية رئيسية", "4 main statistical cards"),
            ("معلومات الاشتراك للأطباء", "Subscription information for doctors"),
            ("مواعيد اليوم والغد للأطباء والمساعدين", "Today's and tomorrow's appointments"),
            ("تنبيهات انتهاء الاشتراك", "Subscription expiration alerts")
        ],
        [
            "📊 4 بطاقات إحصائية: إجمالي الأطباء، المرضى، التشخيصات، تشخيصات اليوم",
            "🎨 كل بطاقة بلون مختلف (Primary, Success, Info, Warning)",
            "📅 بطاقة اشتراك الطبيب (نوع الاشتراك، تاريخ البدء/الانتهاء، الأيام المتبقية)",
            "⚠️ تحذير إذا كان الاشتراك سينتهي خلال 7 أيام",
            "📆 جدولين: مواعيد اليوم ومواعيد الغد مع تفاصيل كل موعد",
            "ℹ️ بطاقة معلومات النظام (نوع المستخدم، التاريخ، الوقت)"
        ],
        LIGHT_BLUE
    )

    # ============ Section: Departments & Specialists ============
    add_title_slide(
        "إدارة الأقسام والتخصصات",
        "Departments & Specialists Management",
        "تنظيم الأقسام الطبية والتخصصات",
        "Organizing Medical Departments and Specialties"
    )

    # ============ Page: Departments Index ============
    add_page_slide(
        "صفحة الأقسام (Departments)",
        "Departments Page",
        [
            ("عرض جميع الأقسام الطبية في جدول", "Display all medical departments in a table"),
            ("زر إضافة قسم جديد", "Add new department button"),
            ("أزرار تعديل وحذف وعرض تفاصيل", "Edit, delete, and view details buttons"),
            ("تنسيق احترافي مع Bootstrap", "Professional formatting with Bootstrap")
        ],
        [
            "📋 جدول (Table) بتنسيق Striped و Hover",
            "🆕 زر \"Add New Department\" باللون الأخضر",
            "📝 أعمدة: الرقم، اسم القسم، الإجراءات",
            "🔧 أزرار الإجراءات: Details (أزرق)، Edit (أصفر)، Delete (أحمر)",
            "🎨 Card مع Card-body لتغليف الجدول",
            "🔍 Responsive Table للعرض على الموبايل"
        ],
        GREEN
    )

    # ============ Page: Specialists Index ============
    add_page_slide(
        "صفحة التخصصات (Specialists)",
        "Specialists Page",
        [
            ("عرض التخصصات مع أقسامها", "Display specialties with their departments"),
            ("Badges ملونة للأقسام", "Colored badges for departments"),
            ("عرض الوصف لكل تخصص", "Description display for each specialty"),
            ("ترقيم تلقائي للصفوف", "Automatic row numbering")
        ],
        [
            "📋 جدول يعرض: الرقم، اسم التخصص، القسم، الوصف، الإجراءات",
            "🏷️ Badge باللون الأزرق لعرض اسم القسم",
            "📝 عمود Description يعرض تفاصيل التخصص",
            "🔗 ربط كل تخصص بقسمه (Foreign Key)",
            "✨ تصميم نظيف ومنظم مع أيقونات Font Awesome",
            "📱 Responsive Design يعمل على جميع الأجهزة"
        ],
        ORANGE
    )

    # ============ Section: Doctors Management ============
    add_title_slide(
        "إدارة الأطباء",
        "Doctors Management",
        "إدارة كاملة لبيانات الأطباء",
        "Complete Doctors Data Management"
    )

    # ============ Page: Doctors Index ============
    add_page_slide(
        "صفحة الأطباء (Doctors)",
        "Doctors Page",
        [
            ("عرض صور الأطباء في الجدول", "Display doctor photos in the table"),
            ("معلومات الاشتراك لكل طبيب", "Subscription information for each doctor"),
            ("حالة النشاط (Active/Not Active)", "Activity status (Active/Not Active)"),
            ("عرض التخصص والقسم", "Display specialty and department"),
            ("معلومات الاتصال الكاملة", "Complete contact information")
        ],
        [
            "🖼️ عمود للصور (50x50 px) مع صورة افتراضية إذا لم توجد",
            "👨‍⚕️ اسم الطبيب + اللقب (Title) تحت الاسم",
            "🏷️ Badge للتخصص (باللون الأزرق) + اسم القسم تحته",
            "📞 رقم الهاتف الأول للطبيب",
            "✅ Badge أخضر (Active) أو رمادي (Not Active)",
            "📅 حالة الاشتراك: Active (أخضر)، Expiring Soon (برتقالي)، Expired (أحمر)، No Subscription (رمادي)",
            "🔧 أزرار: Details، Edit، Delete"
        ],
        LIGHT_BLUE
    )

    # ============ Page: Add/Edit Doctor ============
    add_page_slide(
        "إضافة/تعديل طبيب",
        "Add/Edit Doctor Page",
        [
            ("نموذج شامل لبيانات الطبيب", "Comprehensive form for doctor data"),
            ("رفع صورة الطبيب", "Doctor photo upload"),
            ("اختيار التخصص والقسم", "Specialty and department selection"),
            ("معلومات شخصية ومهنية", "Personal and professional information"),
            ("بيانات تسجيل الدخول للطبيب", "Doctor login credentials")
        ],
        [
            "📋 حقول النموذج: الاسم، اللقب، الرقم المدني، الجنس، العنوان",
            "📞 رقمين للهاتف (Tel1, Tel2)",
            "🏥 قائمة منسدلة للتخصص (Specialist Dropdown)",
            "🖼️ حقل رفع الصورة (File Upload) مع Preview",
            "📝 حقل تفاصيل إضافية (Details Textarea)",
            "👤 معلومات المستخدم: Username، Password (للأطباء القادرين على الدخول)",
            "☑️ Checkbox لـ \"Can Login\" و \"Active\"",
            "💾 أزرار: Save (أخضر)، Back to List (رمادي)"
        ],
        GREEN
    )

    # ============ Section: Patients Management ============
    add_title_slide(
        "إدارة المرضى",
        "Patients Management",
        "إدارة سجلات المرضى",
        "Patient Records Management"
    )

    # ============ Page: Patients Index ============
    add_page_slide(
        "صفحة المرضى (Patients)",
        "Patients Page",
        [
            ("قائمة جميع المرضى", "List of all patients"),
            ("عرض الطبيب المسؤول", "Display assigned doctor"),
            ("معلومات الاتصال", "Contact information"),
            ("الرقم المدني", "Civil ID"),
            ("ترقيم تلقائي", "Automatic numbering")
        ],
        [
            "📋 أعمدة الجدول: الرقم، اسم المريض، الرقم المدني، رقم الهاتف 1، رقم الهاتف 2، الطبيب المسؤول، الإجراءات",
            "👨‍⚕️ عرض اسم الطبيب المسؤول عن كل مريض",
            "📱 عرض رقمي الهاتف",
            "🆔 عرض الرقم المدني (Civil ID)",
            "🔧 أزرار الإجراءات: Details، Edit، Delete",
            "🆕 زر \"Add New Patient\" في أعلى الصفحة"
        ],
        PURPLE
    )

    # ============ Page: Add/Edit Patient ============
    add_page_slide(
        "إضافة/تعديل مريض",
        "Add/Edit Patient Page",
        [
            ("نموذج بيانات المريض", "Patient data form"),
            ("ربط المريض بطبيب", "Assign patient to a doctor"),
            ("معلومات الاتصال الكاملة", "Complete contact information"),
            ("التحقق من البيانات", "Data validation")
        ],
        [
            "📋 حقول: اسم المريض، الرقم المدني، رقم الهاتف 1، رقم الهاتف 2، العنوان",
            "👨‍⚕️ قائمة منسدلة لاختيار الطبيب المسؤول",
            "✅ Required Validation على الحقول المهمة",
            "📝 Textarea للعنوان",
            "💾 زر Save لحفظ البيانات",
            "🔙 زر Back to List للعودة للقائمة"
        ],
        PURPLE
    )

    # ============ Section: Appointments ============
    add_title_slide(
        "إدارة المواعيد",
        "Appointments Management",
        "جدولة ومتابعة المواعيد",
        "Scheduling and Tracking Appointments"
    )

    # ============ Page: Appointments Index ============
    add_page_slide(
        "صفحة المواعيد (Appointments)",
        "Appointments Page",
        [
            ("عرض جميع المواعيد", "Display all appointments"),
            ("حالة الموعد (مجدول، مكتمل، ملغي، لم يحضر)", "Appointment status"),
            ("معلومات المريض والطبيب", "Patient and doctor information"),
            ("التاريخ والوقت", "Date and time"),
            ("السبب والملاحظات", "Reason and notes")
        ],
        [
            "📅 أعمدة: التاريخ، الوقت، المريض، الطبيب، السبب، الحالة، الإجراءات",
            "⏰ عرض الوقت بتنسيق 24 ساعة (HH:mm)",
            "👥 أسماء المريض والطبيب مع أيقونات",
            "🏷️ Badges ملونة للحالة:",
            "  • Scheduled (أزرق) - مجدول",
            "  • Completed (أخضر) - مكتمل",
            "  • Cancelled (أحمر) - ملغي",
            "  • No-Show (برتقالي) - لم يحضر",
            "📝 عمود السبب مع اختصار النص الطويل",
            "🔧 أزرار: Details، Edit، Delete"
        ],
        LIGHT_BLUE
    )

    # ============ Page: Create Appointment ============
    add_page_slide(
        "إضافة موعد جديد",
        "Create Appointment Page",
        [
            ("اختيار المريض والطبيب", "Select patient and doctor"),
            ("تحديد التاريخ والوقت", "Set date and time"),
            ("إدخال السبب والملاحظات", "Enter reason and notes"),
            ("تحديد الحالة", "Set status")
        ],
        [
            "👥 قوائم منسدلة لاختيار المريض والطبيب",
            "📅 Date Picker لاختيار التاريخ",
            "⏰ Time Picker لاختيار الوقت",
            "📝 Textarea للسبب (Reason) - حد أقصى 500 حرف",
            "📝 Textarea للملاحظات (Notes)",
            "🏷️ قائمة منسدلة للحالة (Scheduled/Completed/Cancelled/No-Show)",
            "✅ Validation على جميع الحقول المطلوبة",
            "💾 زر Create Appointment"
        ],
        LIGHT_BLUE
    )

    # ============ Page: Appointments Calendar ============
    add_page_slide(
        "تقويم المواعيد",
        "Appointments Calendar Page",
        [
            ("عرض المواعيد في تقويم شهري", "Display appointments in monthly calendar"),
            ("تمييز الأيام التي بها مواعيد", "Highlight days with appointments"),
            ("إمكانية التنقل بين الأشهر", "Navigate between months"),
            ("عرض تفاصيل المواعيد عند النقر", "Show appointment details on click")
        ],
        [
            "📅 تقويم شهري كامل (Monthly Calendar)",
            "🎨 تلوين الأيام حسب عدد المواعيد",
            "🔍 عرض قائمة المواعيد لليوم المحدد",
            "◀️ ▶️ أزرار Previous/Next Month",
            "📊 عدادات: إجمالي المواعيد، مواعيد اليوم، مواعيد الشهر",
            "🔗 روابط سريعة لإضافة موعد جديد"
        ],
        GREEN
    )

    # ============ Section: Diagnoses ============
    add_title_slide(
        "إدارة التشخيصات الطبية",
        "Medical Diagnoses Management",
        "توثيق التشخيصات الطبية",
        "Documenting Medical Diagnoses"
    )

    # ============ Page: Diagnoses Index ============
    add_page_slide(
        "صفحة التشخيصات (Diagnoses)",
        "Diagnoses Page",
        [
            ("قائمة جميع التشخيصات", "List of all diagnoses"),
            ("معلومات المريض والطبيب", "Patient and doctor information"),
            ("تاريخ التشخيص", "Diagnosis date"),
            ("المرفقات (ملفات PDF)", "Attachments (PDF files)"),
            ("حالة النشاط", "Active status")
        ],
        [
            "📋 أعمدة: التاريخ، المريض، الطبيب، تفاصيل التشخيص، الملف المرفق، الحالة، الإجراءات",
            "👨‍⚕️ أسماء المريض والطبيب مع أيقونات",
            "📅 تنسيق التاريخ (dd MMM yyyy)",
            "📄 رابط لتحميل ملف التشخيص إن وجد",
            "✅ Badge أخضر (Active) أو رمادي (Inactive)",
            "📝 اختصار تفاصيل التشخيص الطويلة",
            "🔧 أزرار: Details، Edit، Delete"
        ],
        ORANGE
    )

    # ============ Page: Add/Edit Diagnosis ============
    add_page_slide(
        "إضافة/تعديل تشخيص",
        "Add/Edit Diagnosis Page",
        [
            ("اختيار المريض والطبيب", "Select patient and doctor"),
            ("تاريخ التشخيص", "Diagnosis date"),
            ("تفاصيل التشخيص", "Diagnosis details"),
            ("رفع ملف PDF", "Upload PDF file"),
            ("تحديد حالة النشاط", "Set active status")
        ],
        [
            "👥 قوائم منسدلة للمريض والطبيب",
            "📅 Date Picker لتاريخ التشخيص",
            "📝 Textarea لتفاصيل التشخيص (كبير ومرن)",
            "📎 File Upload للملف (PDF فقط)",
            "☑️ Checkbox لـ Active Status",
            "✅ Validation على الحقول المطلوبة",
            "💾 زر Save Diagnosis"
        ],
        ORANGE
    )

    # ============ Section: Staff Management ============
    add_title_slide(
        "إدارة الموظفين",
        "Staff Management",
        "إدارة المساعدين والموظفين",
        "Assistants and Staff Management"
    )

    # ============ Page: Doctor Assists Index ============
    add_page_slide(
        "صفحة مساعدي الأطباء",
        "Doctor Assistants Page",
        [
            ("قائمة جميع المساعدين", "List of all assistants"),
            ("الطبيب المسؤول", "Assigned doctor"),
            ("معلومات الاتصال", "Contact information"),
            ("حالة تسجيل الدخول", "Login status"),
            ("حالة النشاط", "Active status")
        ],
        [
            "📋 أعمدة: اسم المساعد، الطبيب، رقم الهاتف 1، رقم الهاتف 2، يمكنه الدخول، نشط، الإجراءات",
            "👨‍⚕️ اسم الطبيب المسؤول مع أيقونة",
            "📱 رقمي الهاتف",
            "🔐 Badge يوضح إذا كان يمكنه تسجيل الدخول (Can Login)",
            "✅ Badge للحالة النشطة (Active/Inactive)",
            "🔧 أزرار الإجراءات"
        ],
        RED_ACCENT
    )

    # ============ Page: My Assistants (For Doctors) ============
    add_page_slide(
        "صفحة مساعديني (للأطباء)",
        "My Assistants Page (For Doctors)",
        [
            ("قائمة مساعدي الطبيب الحالي فقط", "List only current doctor's assistants"),
            ("إضافة مساعد جديد", "Add new assistant"),
            ("إدارة معلومات المساعدين", "Manage assistant information")
        ],
        [
            "👥 عرض مساعدي الطبيب المسجل دخوله فقط",
            "➕ زر \"Add New Assistant\" مخصص للطبيب",
            "📋 نفس مكونات صفحة Doctor Assists ولكن مفلترة",
            "🔒 الطبيب يرى مساعديه فقط",
            "✏️ يمكن للطبيب تعديل بيانات مساعديه"
        ],
        RED_ACCENT
    )

    # ============ Section: Users & Subscriptions ============
    add_title_slide(
        "المستخدمين والاشتراكات",
        "Users & Subscriptions",
        "إدارة المستخدمين واشتراكات الأطباء",
        "Managing Users and Doctor Subscriptions"
    )

    # ============ Page: Users Index ============
    add_page_slide(
        "صفحة المستخدمين (Users)",
        "Users Page",
        [
            ("قائمة جميع مستخدمي النظام", "List of all system users"),
            ("الأدوار والصلاحيات", "Roles and permissions"),
            ("معلومات الاتصال", "Contact information"),
            ("حالة النشاط", "Active status"),
            ("الوظيفة", "Job title")
        ],
        [
            "📋 أعمدة: اسم المستخدم، الاسم الكامل، رقم الهاتف، الوظيفة، الدور، نشط، الإجراءات",
            "🎭 Badge للدور (Role) بألوان مختلفة:",
            "  • Super Admin (أزرق داكن)",
            "  • Admin (أزرق)",
            "  • Doctor (أخضر)",
            "  • Assistant (برتقالي)",
            "  • Receptionist (أصفر)",
            "✅ Badge للحالة النشطة",
            "🔧 أزرار: Details، Edit، Delete"
        ],
        DARK_BLUE
    )

    # ============ Page: Doctor Subscriptions ============
    add_page_slide(
        "صفحة اشتراكات الأطباء",
        "Doctor Subscriptions Page",
        [
            ("إدارة اشتراكات الأطباء", "Manage doctor subscriptions"),
            ("تواريخ البدء والانتهاء", "Start and end dates"),
            ("نوع الاشتراك", "Subscription type"),
            ("الحالة والأيام المتبقية", "Status and remaining days"),
            ("تنبيهات الانتهاء", "Expiration alerts")
        ],
        [
            "📋 أعمدة: الطبيب، نوع الاشتراك، تاريخ البدء، تاريخ الانتهاء، نشط، الحالة، الإجراءات",
            "👨‍⚕️ اسم الطبيب مع تخصصه",
            "📅 تنسيق التواريخ واضح",
            "🏷️ Badges ملونة للحالة:",
            "  • Active (أخضر) - نشط",
            "  • Expiring Soon (برتقالي) - سينتهي قريباً",
            "  • Expired (أحمر) - منتهي",
            "📊 عرض الأيام المتبقية",
            "✅ Checkbox للحالة النشطة"
        ],
        PURPLE
    )

    # ============ Section: Reports ============
    add_title_slide(
        "التقارير والإحصائيات",
        "Reports & Statistics",
        "تقارير شاملة وقابلة للتصدير",
        "Comprehensive and Exportable Reports"
    )

    # ============ Page: Reports Index ============
    add_page_slide(
        "صفحة التقارير الرئيسية",
        "Main Reports Page",
        [
            ("6 أنواع من التقارير", "6 types of reports"),
            ("بطاقات كبيرة وجذابة", "Large attractive cards"),
            ("تصدير Excel", "Excel export"),
            ("فلترة متقدمة", "Advanced filtering")
        ],
        [
            "📊 6 بطاقات تقارير رئيسية:",
            "  1️⃣ تقرير الأطباء (Doctors Report) - أزرق",
            "  2️⃣ تقرير المرضى (Patients Report) - أخضر",
            "  3️⃣ تقرير التشخيصات (Diagnoses Report) - برتقالي",
            "  4️⃣ تقرير الإحصائيات (Statistics Report) - أحمر",
            "  5️⃣ تقرير المساعدين (Assistants Report) - بنفسجي",
            "  6️⃣ تقرير مخصص (Custom Report) - وردي",
            "🎨 كل بطاقة بأيقونة كبيرة وعنوان ووصف",
            "🔗 روابط للذهاب لكل تقرير"
        ],
        GREEN
    )

    # ============ Page: Doctors Report ============
    add_page_slide(
        "تقرير الأطباء",
        "Doctors Report Page",
        [
            ("تقرير شامل لجميع الأطباء", "Comprehensive report of all doctors"),
            ("فلترة حسب التخصص", "Filter by specialty"),
            ("تصدير إلى Excel", "Export to Excel"),
            ("معلومات كاملة لكل طبيب", "Complete information for each doctor")
        ],
        [
            "🔍 قائمة منسدلة لفلترة حسب التخصص",
            "📥 زر \"Export to Excel\" باللون الأخضر",
            "📋 جدول يعرض: الرقم، الاسم، التخصص، القسم، الهاتف، الحالة",
            "📊 إحصائيات: إجمالي الأطباء، النشطين، غير النشطين",
            "🎨 تصميم نظيف مع Badges ملونة",
            "📄 تصدير Excel مع تنسيق احترافي (ClosedXML)"
        ],
        LIGHT_BLUE
    )

    # ============ Page: Patients Report ============
    add_page_slide(
        "تقرير المرضى",
        "Patients Report Page",
        [
            ("تقرير شامل لجميع المرضى", "Comprehensive report of all patients"),
            ("فلترة حسب الطبيب", "Filter by doctor"),
            ("معلومات الاتصال الكاملة", "Complete contact information"),
            ("تصدير Excel", "Excel export")
        ],
        [
            "🔍 قائمة منسدلة لفلترة حسب الطبيب",
            "📥 زر تصدير Excel",
            "📋 أعمدة: الاسم، الرقم المدني، الهاتف، العنوان، الطبيب",
            "👨‍⚕️ عرض اسم الطبيب المسؤول",
            "📊 إحصائيات عدد المرضى",
            "📄 ملف Excel منسق ومنظم"
        ],
        GREEN
    )

    # ============ Page: Diagnoses Report ============
    add_page_slide(
        "تقرير التشخيصات",
        "Diagnoses Report Page",
        [
            ("تقرير تفصيلي للتشخيصات", "Detailed diagnoses report"),
            ("فلترة متعددة (تاريخ، طبيب، مريض)", "Multiple filters (date, doctor, patient)"),
            ("نطاقات التاريخ", "Date ranges"),
            ("تفاصيل التشخيص", "Diagnosis details")
        ],
        [
            "📅 فلتر نطاق التاريخ (From Date - To Date)",
            "👨‍⚕️ قائمة منسدلة للطبيب",
            "👥 قائمة منسدلة للمريض",
            "🔍 زر \"Filter\" لتطبيق الفلاتر",
            "📥 زر \"Export to Excel\"",
            "📋 جدول: التاريخ، المريض، الطبيب، التفاصيل، الملف",
            "📊 إحصائيات: إجمالي التشخيصات، النشطة، غير النشطة"
        ],
        ORANGE
    )

    # ============ Page: Statistics Report ============
    add_page_slide(
        "تقرير الإحصائيات",
        "Statistics Report Page",
        [
            ("إحصائيات شاملة للنظام", "Comprehensive system statistics"),
            ("رسوم بيانية (Charts)", "Charts and graphs"),
            ("بطاقات إحصائية", "Statistical cards"),
            ("تقرير قابل للطباعة", "Printable report")
        ],
        [
            "📊 8 بطاقات إحصائية رئيسية:",
            "  • إجمالي الأطباء، الأطباء النشطين",
            "  • إجمالي المرضى",
            "  • إجمالي التشخيصات، تشخيصات الشهر",
            "  • إجمالي المواعيد، مواعيد اليوم",
            "📈 رسوم بيانية (إن وجدت)",
            "🎨 بطاقات ملونة بألوان مختلفة",
            "🖨️ زر Print Report"
        ],
        RED_ACCENT
    )

    # ============ Section: Profile & Settings ============
    add_title_slide(
        "الملف الشخصي والإعدادات",
        "Profile & Settings",
        "إدارة الملف الشخصي",
        "Managing Personal Profile"
    )

    # ============ Page: Admin Profile ============
    add_page_slide(
        "الملف الشخصي للمدير",
        "Admin Profile Page",
        [
            ("عرض معلومات المدير", "Display admin information"),
            ("تعديل البيانات الشخصية", "Edit personal data"),
            ("تغيير كلمة المرور", "Change password"),
            ("معلومات الحساب", "Account information")
        ],
        [
            "👤 بطاقة معلومات المستخدم:",
            "  • اسم المستخدم (Username)",
            "  • الاسم الكامل (Full Name)",
            "  • رقم الهاتف (Phone)",
            "  • الوظيفة (Job Title)",
            "  • الدور (Role)",
            "✏️ زر \"Edit Profile\" للتعديل",
            "🔐 زر \"Change Password\" لتغيير كلمة المرور",
            "📅 آخر تسجيل دخول (Last Login)"
        ],
        DARK_BLUE
    )

    # ============ Page: Doctor Profile ============
    add_page_slide(
        "الملف الشخصي للطبيب",
        "Doctor Profile Page",
        [
            ("معلومات الطبيب الكاملة", "Complete doctor information"),
            ("الصورة الشخصية", "Profile picture"),
            ("معلومات التخصص", "Specialty information"),
            ("معلومات الاشتراك", "Subscription information")
        ],
        [
            "🖼️ صورة الطبيب الشخصية (كبيرة)",
            "📋 بطاقة معلومات:",
            "  • الاسم واللقب",
            "  • التخصص والقسم",
            "  • الرقم المدني",
            "  • الهاتف والعنوان",
            "📅 بطاقة الاشتراك: نوع، تاريخ البدء/الانتهاء، الأيام المتبقية",
            "✏️ زر تعديل المعلومات",
            "🔐 زر تغيير كلمة المرور"
        ],
        LIGHT_BLUE
    )

    # ============ Page: Change Password ============
    add_page_slide(
        "صفحة تغيير كلمة المرور",
        "Change Password Page",
        [
            ("نموذج آمن لتغيير كلمة المرور", "Secure password change form"),
            ("التحقق من كلمة المرور الحالية", "Verify current password"),
            ("تأكيد كلمة المرور الجديدة", "Confirm new password"),
            ("Validation قوي", "Strong validation")
        ],
        [
            "🔐 3 حقول:",
            "  1. كلمة المرور الحالية (Current Password)",
            "  2. كلمة المرور الجديدة (New Password)",
            "  3. تأكيد كلمة المرور (Confirm Password)",
            "✅ Validation:",
            "  • كلمة المرور الحالية صحيحة",
            "  • كلمة المرور الجديدة مختلفة عن الحالية",
            "  • التأكيد يطابق كلمة المرور الجديدة",
            "💾 زر \"Change Password\"",
            "🔙 زر \"Cancel\" للإلغاء"
        ],
        RED_ACCENT
    )

    # ============ Section: Special Pages ============
    add_title_slide(
        "صفحات خاصة",
        "Special Pages",
        "صفحات إضافية مميزة",
        "Additional Special Pages"
    )

    # ============ Page: Doctor Reports Section ============
    add_page_slide(
        "تقارير الطبيب الشخصية",
        "Doctor Personal Reports",
        [
            ("تقارير خاصة بالطبيب", "Doctor-specific reports"),
            ("مرضاي (My Patients)", "My patients"),
            ("مواعيدي (My Appointments)", "My appointments")
        ],
        [
            "👥 تقرير \"مرضاي\" (My Patients):",
            "  • قائمة مرضى الطبيب فقط",
            "  • إحصائيات عدد المرضى",
            "  • تصدير Excel",
            "📅 تقرير \"مواعيدي\" (My Appointments):",
            "  • مواعيد الطبيب فقط",
            "  • فلترة حسب التاريخ والحالة",
            "  • إحصائيات المواعيد",
            "📊 كل تقرير مفلتر حسب الطبيب المسجل دخوله"
        ],
        GREEN
    )

    # ============ Page: Assistant Reports Section ============
    add_page_slide(
        "تقارير المساعد",
        "Assistant Reports",
        [
            ("تقارير خاصة بالمساعد", "Assistant-specific reports"),
            ("مرضى الطبيب المسؤول", "Assigned doctor's patients"),
            ("المواعيد التي أنشأها المساعد", "Appointments created by assistant")
        ],
        [
            "👨‍⚕️ تقرير مرضى الطبيب (Doctor Patients):",
            "  • مرضى الطبيب المسؤول عن المساعد",
            "  • إحصائيات وتصدير",
            "📅 تقرير المواعيد المُنشأة (My Created Appointments):",
            "  • المواعيد التي أنشأها المساعد فقط",
            "  • معلومات كاملة عن كل موعد",
            "🔒 صلاحيات محدودة حسب الدور"
        ],
        ORANGE
    )

    # ============ Page: Delete Pages ============
    add_page_slide(
        "صفحات الحذف",
        "Delete Pages",
        [
            ("تأكيد الحذف قبل التنفيذ", "Confirm before deletion"),
            ("عرض التفاصيل المراد حذفها", "Display details to be deleted"),
            ("تحذيرات واضحة", "Clear warnings"),
            ("إمكانية الإلغاء", "Ability to cancel")
        ],
        [
            "⚠️ رسالة تحذير كبيرة (Alert Box)",
            "📋 عرض جميع تفاصيل العنصر المراد حذفه",
            "❓ سؤال تأكيد: \"Are you sure you want to delete this?\"",
            "🔴 زر \"Delete\" باللون الأحمر",
            "⚪ زر \"Cancel\" باللون الرمادي",
            "🔒 بعض الصفحات تستخدم Soft Delete (حذف منطقي)",
            "📝 حقل لإدخال سبب الحذف (في صفحة Appointments)"
        ],
        RED_ACCENT
    )

    # ============ Page: Details Pages ============
    add_page_slide(
        "صفحات التفاصيل",
        "Details Pages",
        [
            ("عرض تفصيلي لجميع البيانات", "Detailed display of all data"),
            ("تنسيق منظم ومنسق", "Organized and formatted layout"),
            ("أيقونات توضيحية", "Descriptive icons"),
            ("أزرار إجراءات سريعة", "Quick action buttons")
        ],
        [
            "📄 بطاقة كبيرة (Card) تعرض جميع التفاصيل",
            "🏷️ Labels مع أيقونات Font Awesome",
            "📋 عرض منظم للبيانات (Definition List أو Table)",
            "🖼️ عرض الصور إن وجدت",
            "📎 روابط للملفات المرفقة",
            "🔧 أزرار في الأسفل:",
            "  • Edit (أصفر)",
            "  • Delete (أحمر)",
            "  • Back to List (رمادي)"
        ],
        PURPLE
    )

    # ============ Final Slide: Summary ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ملخص صفحات النظام"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "System Pages Summary"
    p.font.size = Pt(40)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    tf = summary_box.text_frame
    tf.word_wrap = True

    content = [
        "📝 1 صفحة تسجيل دخول (Login Page)",
        "🏠 1 لوحة تحكم رئيسية (Dashboard)",
        "📊 17 مجموعة صفحات إدارية (CRUD Pages)",
        "📈 6 صفحات تقارير (Reports)",
        "👤 3 صفحات ملف شخصي (Profile Pages)",
        "📅 1 صفحة تقويم المواعيد (Calendar)",
        "📱 جميع الصفحات Responsive ومتجاوبة",
        "🎨 تصميم موحد مع Bootstrap 5 و Font Awesome",
        "🔒 نظام صلاحيات كامل لكل صفحة",
        "✅ نماذج بـ Validation كاملة"
    ]

    for item in content:
        p = tf.paragraphs[0] if item == content[0] else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = thanks_box.text_frame

    p = tf.paragraphs[0]
    p.text = "شكراً"
    p.font.size = Pt(70)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "نظام إدارة العيادات - Clinic Management System"
    p.font.size = Pt(22)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "/home/user/ClinicManagementSystem/ClinicManagementSystem_Pages_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        output_file = create_detailed_presentation()
        print(f"\n✓ Detailed PowerPoint presentation created successfully!")
        print(f"✓ File: {output_file}")
        print(f"✓ Total slides: 40+")
        print(f"✓ Language: Bilingual (Arabic/English)")
        print(f"✓ Content: All website pages with detailed descriptions")
    except Exception as e:
        print(f"✗ Error creating presentation: {str(e)}")
        raise
